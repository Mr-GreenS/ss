import time
import telebot

import os
import docx2txt
from pptx import Presentation
import pdfplumber

def findFiles(strings, dir, subDirs, fileContent, fileExtensions):
    # Finds all the files in 'dir' that contain one string from    'strings'. 
    # Additional parameters:
    # 'subDirs': True/False : Look in sub-directories of your folder
    # 'fileContent': True/False :Also look for the strings in the file     content of every file
    # 'fileExtensions': True/False : Look for a specific file extension -> 'fileContent' is ignored
    filesInDir = []
    foundFiles = []
    filesFound = 0

    if not subDirs:
        for filename in os.listdir(dir):
            if os.path.isfile(os.path.join(dir, filename).replace("\\", "/")):
                filesInDir.append(os.path.join(dir, filename).replace("\\", "/"))
    else:
        for root, subdirs, files in os.walk(dir):
            for f in files:
                if not os.path.isdir(os.path.join(root, f).replace("\\", "/")):
                    filesInDir.append(os.path.join(root, f).replace("\\", "/"))
    print(filesInDir)
    # Find files that contain the keyword
    if filesInDir:
        for file in filesInDir:
            print("Current file: "+file)
            # Define what is to be searched in
            filename, extension = os.path.splitext(file)
            if fileExtensions:
                fileText = extension
            else:
                fileText = os.path.basename(filename).lower()
                if fileContent:
                    fileText +=  getFileContent(file).lower()
            # Check for translations
            for string in strings:
                print(string)
                if string in fileText:
                    foundFiles.append(file)
                    filesFound += 1
                    break
    return foundFiles

def getFileContent(filename):
    '''Returns the content of a file of a supported type (list: supportedTypes)'''
    if filename.partition(".")[2] in supportedTypes:
        if filename.endswith(".pdf"):
            content = ""
            with pdfplumber.open(filename) as pdf:
                for x in range(0, len(pdf.pages)):
                    page = pdf.pages[x]
                    content = content + page.extract_text()
            return content
        elif filename.endswith(".txt"):
            with open(filename, 'r') as f:
                content = ""
                lines = f.readlines()
                for x in lines:
                    content = content + x
            f.close()
            return content
        elif filename.endswith(".docx"):
            content = docx2txt.process(filename)
            return content
        elif filename.endswith(".pptx"):
            content = ""
            prs = Presentation(filename)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content = content+shape.text
            return content
    else:
        return ""

supportedTypes = ["txt", "docx", "pdf", "pptx"]
# print(findFiles(strings=["buch"], dir="C:/Users/User/Desktop/",  subDirs=True, fileContent=True, fileExtensions=False))
i_pass= 0
SPASS= "Aimee"

bot = telebot.TeleBot("6065910250:AAEJ0P0eFCP3rmr5S6WLw0452LuJxDgKlt4", parse_mode=None) # You can set parse_mode by default. HTML or MARKDOWN
@bot.message_handler(commands=['start'])
def send_welcome(message):
    bot.send_message(message.chat.id, "\n\n=======AI-M101101 FOREVER=======\n\n")
    bot.send_message(message.chat.id, "\n\nВведите пароль\n\n")

@bot.message_handler(commands=['help'])
def send_help(message):
    scomm="\n\nЗапрос Команда\n\n"
    scomm+="Примеры:"
    scomm+="Dir test\nFile Сертификат\nFile >./test/x/Aimee0.png "
    bot.send_message(message.chat.id, scomm)

def to_dir(message, command):
    print("DIR")
    print(command)
    s=s1=''
    if command[2] != '':
       print(command[2])
       list=os.listdir(command[2])
    else:
        list=os.listdir();
    print(list)
    for s in list:
        print(s)
        s1=s1+ '\n' + s
    if s1 != '':
            return s1
    else: return "Директория пустая"

   # распечатать все файлы и папки рекурсивно
            #for dirpath, dirnames, filenames in os.walk("."):
                # перебрать каталоги
                #for dirname in dirnames:
                #	bot.send_message(message.chat.id,"Каталог:" + os.path.join(dirpath, dirname))
                # перебрать файлы
                #for filename in filenames:
                #	bot.send_message(message.chat.id,"Файл:" + os.path.join(dirpath, filename))
       
            #	os.mkdir(message.text)
            #bot.send_message(message.chat.id, "Создана директория:" + message.text)
            #time.sleep(20)
            #bot.stop_bot()
def to_file(message, command):
    print("FILE")
    # sendDocument
    doc = open('result.txt', 'w')
   
    s=''
    if command[2] != "":
        list=findFiles(strings=[command[2]], dir=".",  subDirs=True, fileContent=True, fileExtensions=False)
        for s in list:
           doc.write(s+'\n')
        if s != '':
               print(s)
               doc.close()
               with open("result.txt","rb") as misc:
                   try:
                       bot.send_document(message.chat.id, misc, caption="Вот что нашлось")
                   except:
                       print("Exception send_document")
                   s="Отправлен файл результата, спасибо"
               os.remove('result.txt')
               return s
        else:
               return "File not found"
    else:
     print(command)
     return "File not found"

def findJira(command, num):
    try:
        n=int(command)
        print('find by num ' + str(n))
    except:
        print('find by str '+ command)

    list =[]
    return "Tickets not found"

def to_jira(message, command):
    print("JIRA")
    s=s1=''
    if command[2] != "":
        list=findJira(command[2], 200)
        for s in list:
           s1=s1 + '\n' + s
           if s1 != '':
               return s1
           else:
               return "Tickets not found"
    else:
     print(command)
     return "Ticket not found"

def to_aimee(message, command):
    print("Эми- OpenAI bot")
    if True:
        print("Эта часть не дописана")
        bot.send_photo(message.chat.id, photo=open('./test/x/Aimee1.png', 'rb'))

    return "Эми в отпуске"

@bot.message_handler(content_types=['text'])
def send_echo(message):
    #bot.reply_to(message, message.text)
    global i_pass
    scomm=str(message.text)
    command=scomm.partition(' ')
    print(command)
    try:
        if message.text == SPASS :
            os.chdir("D:\\AIMEE")
            bot.send_message(message.chat.id, "Привет. Эми активирована.\nТекущий путь: " + os.getcwd())
#           bot.send_photo(message.chat.id, photo=open('./test/x/Aimee0.png', 'rb'))
            i_pass = 1
        elif i_pass == 1:
            if command[0] == "DIR" or command[0] =="Dir":#обрабатываем команду директории
                scomm = to_dir(message,command)
            elif command[0] == "FILE" or command[0] == "File":#обрабатываем команду файла
                print("FILE SECTION")
                if command[2]!='':#отправить указанный файл
                    if command[2][0]=='>' and len(command[2])>1:
                            print("send file:" + str(command[2][1:]))
                            with open(command[2][1:],"rb") as misc:
                                if misc != None:
                                    bot.send_document(message.chat.id, misc, caption="Вот ваш файл")
                                scomm="Отправлен файл результата, спасибо"
                else:#найти файл, отправить результат
                        scomm = to_file(message,command)
                
            elif command[0] == "JIRA" or command[0] =="Jira":#обрабатываем команду жира
                scomm = to_jira(message,command)
            elif command[0] == "Эми" or command[0] =="ЭМИ":#обрабатываем команду ai
                scomm = to_aimee(message,command)
        else:
            scomm = "Запрос не распознан"

        if scomm != '' and scomm != None:
            bot.send_message(message.chat.id,scomm)
        else:
            bot.send_message(message.chat.id, "Запрос не распознан")
    except:
        print("Unpredictable error")

print("Bot polling")
bot.polling(none_stop = True)
time.sleep(2.5)




